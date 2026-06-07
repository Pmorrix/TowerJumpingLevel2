from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


SRC = Path(r"C:\Users\Phillips\Documents\IA_Generativa_ChatGPT_Videojuegos_Clase_v2_Master_UMA.pptx")
OUT = Path(r"C:\Users\Phillips\Documents\LLM_ChatGPT_Go_Unity_Clase1_Master_UMA.pptx")


SLIDES = {
    1: [
        "LLM y\nChatGPT Go",
        "Uso profesional en desarrollo de videojuegos con Unity 3D y C#",
        "Clase 1 - 5 horas - fundamentos, prompts, riesgos y revisión técnica",
        "5h",
        "clase\nteórica-práctica",
        "Máster en Creación de Videojuegos - Universidad de Málaga",
    ],
    2: [
        "ENFOQUE",
        "Enfoque",
        "02",
        "Objetivo de la clase",
        "Qué debe dominar el alumnado antes de pasar al taller práctico",
        "Comprender qué es un LLM y cómo trabaja con tokens, contexto y predicción.",
        "Usar ChatGPT Go como asistente técnico para Unity, no como autoridad automática.",
        "Formular prompts útiles con contexto, restricciones y formato verificable.",
        "Detectar alucinaciones, errores plausibles y riesgos específicos de Unity.",
        "Preparar una forma de trabajo segura para la clase práctica posterior.",
        "Idea central",
        "ChatGPT acelera el aprendizaje y el prototipado cuando el problema está bien formulado y la respuesta se verifica.",
    ],
    3: [
        "PLAN DE CLASE",
        "Plan de clase",
        "03",
        "Ruta de 5 horas",
        "Estructura cerrada de la primera sesión",
        "01",
        "Apertura",
        "20 min",
        "02",
        "ChatGPT Go",
        "25 min",
        "03",
        "LLM",
        "40 min",
        "04",
        "Transformers",
        "35 min",
        "05",
        "Descanso",
        "10 min",
        "06",
        "Videojuegos",
        "35 min",
        "07",
        "Prompts Unity",
        "45 min",
        "08",
        "Riesgos",
        "40 min",
        "09",
        "Revisión código",
        "30 min",
        "10",
        "Cierre",
        "10 min",
    ],
    4: [
        "FUNDAMENTOS",
        "Fundamentos",
        "04",
        "IA generativa en contexto",
        "Diferencia práctica para un desarrollador de videojuegos",
        "IA clásica",
        "Sistemas orientados a decidir, clasificar, predecir o controlar a partir de reglas, datos o modelos entrenados.\n\nEjemplos: pathfinding, FSM, behavior trees, detección, scoring.",
        "IA generativa",
        "Sistemas orientados a producir contenido nuevo: texto, código, imágenes, audio, ideas o variantes.\n\nEjemplos: diálogos, misiones, documentación, prototipos de scripts, QA.",
    ],
    5: [
        "CHATGPT GO",
        "ChatGPT Go",
        "05",
        "La herramienta que usaremos",
        "Más margen que el plan gratuito, pero no una garantía de verdad",
        "Permite trabajar con conversaciones más largas y más uso que una cuenta gratuita.",
        "Puede incluir carga de archivos, imágenes, análisis y herramientas según disponibilidad.",
        "No es lo mismo que la API: aquí trabajamos desde la interfaz de ChatGPT.",
        "Sus límites pueden variar por cuenta, región, carga del sistema o cambios del servicio.",
        "Para clase",
        "Lo importante no es la suscripción, sino el método: pedir con precisión, revisar y probar.",
    ],
    6: [
        "FUNDAMENTOS",
        "Fundamentos",
        "06",
        "Tokens: el texto se trocea",
        "El modelo no recibe frases como las vemos nosotros",
        "Texto",
        "La ruta natural",
        "Tokens",
        '["La", " ruta", " natural"]',
        "IDs",
        "[4579, 59781, 6247]",
        "Criterio",
        "Los números son índices de vocabulario. No significan por sí solos: el significado aparece en el contexto y en los vectores internos.",
    ],
    7: [
        "FUNDAMENTOS",
        "Fundamentos",
        "07",
        "Embeddings y vectores",
        "Representaciones numéricas de alta dimensión",
        "Token -> ID -> vector",
        "Un fragmento de texto pasa a un número de vocabulario y después a una lista larga de valores.",
        "Idea clave",
        "No son coordenadas físicas: son coordenadas matemáticas para que el modelo opere con lenguaje.",
    ],
    8: [
        "CONTEXTO",
        "Contexto",
        "08",
        "Qué contexto necesita ChatGPT",
        "En Unity, el código aislado rara vez cuenta toda la historia",
        "Proyecto",
        "género, objetivo, nivel técnico",
        "Escena",
        "GameObjects, prefabs, tags, layers",
        "Scripts",
        "responsabilidades y dependencias",
        "Inspector",
        "SerializeField, UI, botones, referencias",
        "Restricciones",
        "versión, paquetes, físicas, Time.timeScale",
        "Principio operativo: cuanto mejor describes el entorno, menos tiene que inventar.",
    ],
    9: [
        "TRANSFORMERS",
        "Transformers",
        "09",
        "Attention Is All You Need",
        "Por qué el Transformer cambió la IA generativa moderna",
        "Antes",
        "Muchos modelos procesaban secuencias de forma más lineal o recurrente.",
        "Transformer",
        "Permite relacionar muchas partes del contexto mediante atención.",
        "Ventaja",
        "Entrenamiento más paralelizable y mejor manejo de dependencias largas.",
        "Para nosotros: ayuda a entender por qué ChatGPT relaciona instrucciones, código y restricciones dentro de una conversación.",
    ],
    10: [
        "ATENCIÓN",
        "Atención",
        "10",
        "Atención no es comprensión perfecta",
        "Es una forma de ponderar qué partes importan para responder",
        "Ejemplo",
        "La cámara siguió al jugador porque estaba demasiado cerca.",
        "Pregunta",
        "¿Qué estaba demasiado cerca: la cámara o el jugador?",
        "En código",
        "Relaciona scoreText, TextMeshProUGUI, Inspector y posibles NullReferenceException.",
        "Límite: relacionar partes del contexto no equivale a verificar que la respuesta sea correcta.",
    ],
    11: [
        "VIDEOJUEGOS",
        "Videojuegos",
        "11",
        "Dónde aporta valor",
        "No solo código: también diseño, producción y validación",
        "Diseño: mecánicas, loops, economía, niveles.",
        "Programación: pseudocódigo, scripts base, explicación de errores.",
        "Narrativa: diálogos, NPCs, misiones, lore.",
        "Producción: tareas, mini GDD, roadmap, criterios de aceptación.",
        "QA: casos de prueba, edge cases y checklist.",
        "Principio",
        "Cuanto más concreto sea el problema, más útil será la respuesta.",
    ],
    12: [
        "UNITY",
        "Unity",
        "12",
        "El contexto de Unity importa",
        "La escena también es parte del problema técnico",
        "Inspector y escena",
        "SerializeField, referencias UI, prefabs, tags, layers, botones y GameObjects.",
        "Flujo técnico",
        "Rigidbody, Collider, Trigger, eventos, escenas, FixedUpdate y Time.timeScale.",
    ],
    13: [
        "BLOQUE PRÁCTICO",
        "Prompting",
        "De pedir respuestas a especificar problemas técnicos",
        "13",
    ],
    14: [
        "PROMPTS",
        "Prompts",
        "14",
        "Anatomía de un buen prompt",
        "Una especificación mínima para trabajar con ChatGPT",
        "Rol",
        "Quién debe actuar",
        "Contexto",
        "Proyecto, escena y nivel técnico",
        "Objetivo",
        "Qué resultado se necesita",
        "Restricciones",
        "Qué debe evitar o respetar",
        "Formato",
        "Cómo debe devolver la respuesta",
        "Criterio",
        "Cómo sabremos si sirve",
    ],
    15: [
        "PROMPTS",
        "Prompts",
        "15",
        "Prompt malo vs prompt útil",
        "Comparar para que la mejora sea evidente",
        "Prompt débil",
        '"Hazme un enemigo."\n\nProblemas:\n- no hay tipo de juego\n- no hay escena\n- no hay física\n- no hay restricciones\n- no hay criterio de prueba',
        "Prompt útil",
        '"Actúa como desarrollador Unity C# senior. Necesito un enemigo simple: un cubo que persiga a una cápsula Player a menos de 8 unidades. Usa Vector3.MoveTowards, no NavMesh, y dime qué revisar en Inspector."',
    ],
    16: [
        "PROMPTS",
        "Prompts",
        "16",
        "La respuesta inicial no es el final",
        "La utilidad aparece en la iteración",
        "Pedir alternativas",
        "Dame 3 enfoques más simples.",
        "Pedir reducción",
        "Elimina lo que no sea imprescindible.",
        "Pedir comparación",
        "Compara coste, riesgo y valor jugable.",
        "Pedir revisión",
        "Busca bugs, supuestos ocultos y pruebas necesarias.",
    ],
    17: [
        "REVISIÓN",
        "Revisión",
        "17",
        "ChatGPT como revisor",
        "Útil cuando se le pide mirar riesgos concretos",
        "Primero",
        "Pegar error completo o script completo.",
        "Después",
        "Describir GameObject, componentes y referencias del Inspector.",
        "Pedir",
        "Causa probable, comprobación, cambio mínimo y pruebas.",
        "Evitar",
        'Preguntas vagas como "no funciona" sin escena ni consola.',
        "Regla docente: revisar respuestas con criterio técnico antes de integrarlas.",
    ],
    18: [
        "CHECKLIST",
        "Checklist",
        "18",
        "Antes de aceptar código generado",
        "Una revisión mínima para Unity 3D",
        "OK",
        "¿Compila?",
        "OK",
        "¿Es Unity 3D y no 2D?",
        "OK",
        "¿Respeta SerializeField e Inspector?",
        "OK",
        "¿Usa componentes que existen?",
        "OK",
        "¿Puede producir NullReferenceException?",
        "OK",
        "¿Afecta Time.timeScale, escenas o prefabs?",
    ],
    19: [
        "RIESGOS",
        "Riesgos",
        "La IA puede fallar de forma convincente",
        "19",
    ],
    20: [
        "RIESGOS",
        "Riesgos",
        "20",
        "Alucinaciones en Unity",
        "Errores típicos que parecen soluciones reales",
        "APIs inventadas o paquetes que no existen en la versión del proyecto.",
        "Mezcla 2D/3D: Rigidbody2D en un proyecto 3D o triggers mal planteados.",
        "Complejidad extra: managers, singletons o sistemas no pedidos.",
        "Inspector roto: cambiar nombres públicos o borrar SerializeField usados por escena.",
        "Tiempo y flujo: Time.timeScale, eventos sin desuscribir o escenas sin restaurar estado.",
        "Pregunta clave",
        "¿Qué puede romperse si copio esta respuesta sin revisar?",
    ],
    21: [
        "USO RESPONSABLE",
        "Uso responsable",
        "21",
        "Privacidad y autoría",
        "Usar IA en profesional exige límites claros",
        "No subir código propietario, datos personales o material confidencial sin permiso.",
        "No copiar código generado sin entenderlo ni probarlo.",
        "Citar o declarar el uso de IA si el centro o el proyecto lo exige.",
        "Verificar licencias de assets, imágenes o textos externos.",
        "Mantener trazabilidad de decisiones técnicas importantes.",
        "Mensaje clave",
        "La responsabilidad del resultado sigue siendo del desarrollador.",
        "Buenas prácticas",
    ],
    22: [
        "CHATGPT GO",
        "ChatGPT Go",
        "22",
        "Archivos, imágenes y análisis",
        "Más modos de dar contexto al modelo",
        "Script .cs",
        "Subirlo para explicación, revisión o depuración.",
        "Captura de consola",
        "Pedir causa probable y comprobaciones.",
        "Captura del Inspector",
        "Detectar referencias vacías o componentes faltantes.",
        "Documento de diseño",
        "Extraer requisitos, riesgos y tareas.",
        "Regla",
        "No pedir código hasta haber identificado el problema y sus restricciones.",
    ],
    23: [
        "PRÁCTICA",
        "Práctica",
        "23",
        "Ejercicio de cierre",
        "Aplicar el método sin crear todavía el minijuego",
        "Tarea",
        "Elegir un caso: mecánica, enemigo, UI, bug, sistema de vidas o checklist QA.",
        "Condición",
        "Realizar 3 iteraciones: prompt inicial, mejora con restricciones y revisión crítica.",
        "Entrega oral",
        "Explicar qué cambió entre la primera y la última respuesta, y por qué la versión final es más verificable.",
    ],
    24: [
        "CIERRE",
        "Cierre",
        "24",
        "Mensajes clave",
        "Qué debe quedar claro al final de la sesión",
        "Un LLM trabaja con tokens, contexto y predicción.",
        "La atención relaciona partes del contexto, pero no garantiza verdad.",
        "ChatGPT Go da más margen, no elimina errores.",
        "En Unity, el Inspector y la escena son parte del problema.",
        "La práctica profesional es pedir, revisar, probar y documentar.",
    ],
    25: [
        "Siguiente sesión",
        "Taller práctico:\nminijuego en Unity",
        "Aplicaremos el método con cápsulas, cubos, esferas y planos: movimiento, coleccionables, obstáculos, UI, victoria, derrota y reinicio.",
        "Máster en Creación de Videojuegos - UMA",
    ],
}


def main() -> None:
    prs = Presentation(str(SRC))
    for idx, slide in enumerate(prs.slides, 1):
        replacements = SLIDES[idx]
        text_shapes = [
            shape
            for shape in slide.shapes
            if hasattr(shape, "text_frame")
            and shape.text_frame is not None
            and shape.text.strip()
        ]
        for position, shape in enumerate(text_shapes):
            shape.text = replacements[position] if position < len(replacements) else ""
    rebuild_tokens_slide(prs.slides[5])
    rebuild_prompt_anatomy_slide(prs.slides[13])
    prs.save(str(OUT))
    print(OUT)
    print(f"slides={len(prs.slides)}")


def clear_text(slide) -> None:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            shape.text = ""


def add_text(slide, x, y, w, h, value, size=18, bold=False, color=(20, 20, 20), align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_card(slide, x, y, w, h, title, body, fill=(249, 232, 236), stroke=(230, 20, 45)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(*stroke)
    shape.line.width = Pt(1)
    add_text(slide, x + 0.14, y + 0.12, w - 0.28, 0.25, title, size=13, bold=True, color=(50, 50, 50))
    add_text(slide, x + 0.14, y + 0.44, w - 0.28, h - 0.54, body, size=12, color=(35, 35, 35))


def add_footer(slide, section, number):
    add_text(slide, 0.42, 7.03, 2.0, 0.18, section, size=8, color=(40, 40, 40))
    add_text(slide, 12.05, 7.03, 0.35, 0.18, number, size=8, color=(40, 40, 40))


def rebuild_tokens_slide(slide) -> None:
    clear_text(slide)
    add_text(slide, 0.42, 0.16, 2.3, 0.24, "FUNDAMENTOS", size=8, color=(20, 20, 20))
    add_text(slide, 0.62, 0.82, 8.0, 0.45, "Tokens: el texto se trocea", size=18, bold=True)
    add_text(slide, 0.62, 1.36, 8.0, 0.35, "El modelo no recibe frases como las vemos nosotros", size=13)
    add_card(slide, 0.62, 2.15, 2.2, 1.05, "Texto", "La ruta natural", fill=(252, 238, 241))
    add_text(slide, 2.96, 2.47, 0.35, 0.35, ">", size=28, bold=True, color=(230, 20, 45))
    add_card(slide, 3.36, 2.15, 2.5, 1.05, "Tokens", '["La", " ruta",\n" natural"]', fill=(245, 247, 250), stroke=(200, 205, 214))
    add_text(slide, 6.02, 2.47, 0.35, 0.35, ">", size=28, bold=True, color=(230, 20, 45))
    add_card(slide, 6.42, 2.15, 2.5, 1.05, "IDs", "[4579, 59781,\n6247]", fill=(245, 247, 250), stroke=(200, 205, 214))
    add_card(slide, 9.42, 1.95, 2.65, 1.55, "Criterio", "Los números son índices de vocabulario. El significado aparece en el contexto y en los vectores internos.", fill=(245, 247, 250), stroke=(200, 205, 214))
    add_text(slide, 0.62, 4.5, 10.8, 0.5, "Para Unity: si el contexto no incluye escena, componentes e Inspector, el modelo completa huecos con supuestos.", size=14, bold=True, color=(60, 60, 60))
    add_footer(slide, "Fundamentos", "06")


def rebuild_prompt_anatomy_slide(slide) -> None:
    clear_text(slide)
    add_text(slide, 0.42, 0.16, 2.3, 0.24, "PROMPTS", size=8, color=(20, 20, 20))
    add_text(slide, 0.62, 0.82, 8.0, 0.45, "Anatomía de un buen prompt", size=18, bold=True)
    add_text(slide, 0.62, 1.36, 8.6, 0.35, "Una especificación mínima para trabajar con ChatGPT", size=13)
    cards = [
        ("Rol", "Quién debe actuar"),
        ("Contexto", "Proyecto, escena y nivel técnico"),
        ("Objetivo", "Qué resultado se necesita"),
        ("Restricciones", "Qué debe evitar o respetar"),
        ("Formato", "Cómo debe devolver la respuesta"),
        ("Criterio", "Cómo sabremos si sirve"),
    ]
    for i, (title, body) in enumerate(cards):
        col = i % 3
        row = i // 3
        add_card(slide, 0.62 + col * 3.95, 2.05 + row * 1.55, 3.25, 1.05, title, body, fill=(245, 247, 250), stroke=(200, 205, 214))
    add_text(slide, 0.62, 5.45, 10.7, 0.45, "Prompt útil = instrucciones claras + contexto suficiente + salida evaluable.", size=16, bold=True, color=(230, 20, 45))
    add_footer(slide, "Prompts", "14")


if __name__ == "__main__":
    main()
